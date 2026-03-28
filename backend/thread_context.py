"""Shared parsing of AI SDK / assistant-ui messages: attachments, text extraction, OpenAI blocks."""

from __future__ import annotations

import base64
import io
import logging
from pathlib import Path
from urllib.parse import unquote_to_bytes

from pypdf import PdfReader

logger = logging.getLogger(__name__)

MAX_EXTRACTED_FILE_CHARS = 80_000


def parse_data_url(url: str) -> tuple[str, bytes] | None:
    """Parse a data: URL into (mime_type, raw_bytes). Returns None if invalid."""
    if not isinstance(url, str) or not url.startswith("data:"):
        return None
    comma = url.find(",")
    if comma == -1:
        return None
    meta = url[5:comma]
    payload = url[comma + 1 :]
    if ";base64" in meta:
        mime = meta.split(";base64", 1)[0].strip() or "application/octet-stream"
        try:
            raw = base64.b64decode(payload)
        except Exception:
            return None
    else:
        mime = (meta.split(";", 1)[0].strip() if meta else "") or "text/plain"
        raw = unquote_to_bytes(payload)
    return mime, raw


def extract_pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    chunks: list[str] = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    return "\n\n".join(chunks)


def extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
    return "\n".join(lines)


def iter_message_parts(m: dict) -> list[dict]:
    parts = m.get("parts")
    if parts and isinstance(parts, list):
        return [p for p in parts if isinstance(p, dict)]
    content = m.get("content")
    if isinstance(content, list):
        return [p for p in content if isinstance(p, dict)]
    return []


def file_part_to_openai_blocks(part: dict) -> list[dict]:
    """Turn one AI SDK file part (data URL or https URL) into OpenAI content blocks."""
    url = part.get("url") or ""
    media = (part.get("mediaType") or "").lower()
    filename = part.get("filename") or "attachment"

    if url.startswith("http://") or url.startswith("https://"):
        if media.startswith("image/"):
            return [{"type": "image_url", "image_url": {"url": url}}]
        return [
            {
                "type": "text",
                "text": (
                    f"[Remote file {filename} — only image URLs are loaded directly; "
                    "download and re-attach or paste text.]"
                ),
            }
        ]

    parsed = parse_data_url(url)
    if not parsed:
        return [{"type": "text", "text": f"[Could not decode attachment: {filename}]"}]

    mime, raw = parsed
    mime_l = mime.lower().split(";")[0].strip()

    ext = Path(filename).suffix.lower()
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
    is_image = (
        media.startswith("image/")
        or mime_l.startswith("image/")
        or ext in image_exts
    )
    if is_image:
        if url.startswith("data:"):
            return [{"type": "image_url", "image_url": {"url": url}}]
        b64 = base64.b64encode(raw).decode()
        type_mime = media.split(";")[0].strip() if media.startswith("image/") else mime_l
        return [{"type": "image_url", "image_url": {"url": f"data:{type_mime};base64,{b64}"}}]

    if mime_l == "application/pdf" or filename.lower().endswith(".pdf"):
        text = extract_pdf_text(raw)
        if not text.strip():
            text = "(No extractable text found in this PDF.)"
        text = text[:MAX_EXTRACTED_FILE_CHARS]
        return [
            {
                "type": "text",
                "text": f"--- {filename} (PDF text) ---\n{text}",
            }
        ]

    pptx_mimes = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    if mime_l in pptx_mimes or filename.lower().endswith(".pptx"):
        try:
            text = extract_pptx_text(raw)
        except Exception as e:
            logger.warning("pptx extract failed for %s: %s", filename, e)
            text = ""
        if not text.strip():
            text = "(No extractable text found in this presentation.)"
        text = text[:MAX_EXTRACTED_FILE_CHARS]
        return [
            {
                "type": "text",
                "text": f"--- {filename} (slide text) ---\n{text}",
            }
        ]

    if mime_l.startswith("text/") or mime_l in (
        "application/json",
        "application/xml",
    ):
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("utf-8", errors="replace")
        text = text[:MAX_EXTRACTED_FILE_CHARS]
        return [{"type": "text", "text": f"--- {filename} ---\n{text}"}]

    return [
        {
            "type": "text",
            "text": (
                f"[Attached {filename} ({mime_l}) — unsupported type; "
                "try PDF, PPTX, plain text, or an image.]"
            ),
        }
    ]


def extract_text_from_message(m: dict) -> str:
    """Extract text content from an AI SDK message (text parts only)."""
    parts = m.get("parts")
    if parts and isinstance(parts, list):
        texts = []
        for part in parts:
            if isinstance(part, dict) and part.get("type") == "text":
                texts.append(part.get("text", ""))
        if texts:
            return " ".join(texts)

    content = m.get("content", "")

    if isinstance(content, list):
        texts = []
        for part in content:
            if isinstance(part, dict):
                texts.append(part.get("text", ""))
        result = " ".join(texts)
        if result.strip():
            return result

    if isinstance(content, str) and content.strip():
        return content

    return str(content)


def collapse_adjacent_text_blocks(blocks: list[dict]) -> list[dict]:
    out: list[dict] = []
    buf: list[str] = []
    for b in blocks:
        if isinstance(b, dict) and b.get("type") == "text":
            buf.append(b.get("text", "") or "")
        else:
            if buf:
                out.append({"type": "text", "text": "".join(buf)})
                buf = []
            out.append(b)
    if buf:
        out.append({"type": "text", "text": "".join(buf)})
    return out


def compile_messages_for_slide_generation(messages: list) -> str | list:
    """Build one user payload for slide generation: plain string, or multimodal content list."""
    blocks: list[dict] = [
        {
            "type": "text",
            "text": (
                "Use the full conversation below and all attached materials (including images) "
                "to create the presentation. Honor facts, numbers, and structure from "
                "attachments.\n\n"
            ),
        }
    ]

    for m in messages:
        if not isinstance(m, dict):
            continue
        role = m.get("role", "user")
        if role not in ("user", "assistant"):
            continue
        label = "USER" if role == "user" else "ASSISTANT"
        blocks.append({"type": "text", "text": f"\n### {label}\n"})

        parts = iter_message_parts(m)
        if not parts:
            t = extract_text_from_message(m).strip()
            if t:
                blocks.append({"type": "text", "text": t + "\n"})
            continue

        for part in parts:
            pt = part.get("type")
            if pt == "text":
                t = (part.get("text") or "").strip()
                if t:
                    blocks.append({"type": "text", "text": t + "\n"})
            elif pt == "file":
                blocks.extend(file_part_to_openai_blocks(part))

    collapsed = collapse_adjacent_text_blocks(blocks)
    has_image = any(
        isinstance(b, dict) and b.get("type") == "image_url" for b in collapsed
    )
    if not has_image:
        texts = [
            (b.get("text") or "")
            for b in collapsed
            if isinstance(b, dict) and b.get("type") == "text"
        ]
        full = "".join(texts).strip()
        return full or "Generate a presentation"

    return collapsed


def compile_messages_for_live_qa(messages: list) -> str | list:
    """Build knowledge-base payload for live Q&A (conversation + attachments)."""
    blocks: list[dict] = [
        {
            "type": "text",
            "text": (
                "Knowledge base for a live presentation Q&A. Use ONLY facts from the "
                "conversation and attached materials below. If something is not covered, "
                "say clearly that the materials do not contain enough information—do "
                "not invent specifics.\n\n"
            ),
        }
    ]

    for m in messages:
        if not isinstance(m, dict):
            continue
        role = m.get("role", "user")
        if role not in ("user", "assistant"):
            continue
        label = "USER" if role == "user" else "ASSISTANT"
        blocks.append({"type": "text", "text": f"\n### {label}\n"})

        parts = iter_message_parts(m)
        if not parts:
            t = extract_text_from_message(m).strip()
            if t:
                blocks.append({"type": "text", "text": t + "\n"})
            continue

        for part in parts:
            pt = part.get("type")
            if pt == "text":
                t = (part.get("text") or "").strip()
                if t:
                    blocks.append({"type": "text", "text": t + "\n"})
            elif pt == "file":
                blocks.extend(file_part_to_openai_blocks(part))

    collapsed = collapse_adjacent_text_blocks(blocks)
    has_image = any(
        isinstance(b, dict) and b.get("type") == "image_url" for b in collapsed
    )
    if not has_image:
        texts = [
            (b.get("text") or "")
            for b in collapsed
            if isinstance(b, dict) and b.get("type") == "text"
        ]
        full = "".join(texts).strip()
        return full or "(No chat context was provided.)"

    return collapsed


def build_user_openai_message(m: dict) -> dict:
    """Build a user message for OpenAI (string or multimodal content array)."""
    blocks: list[dict] = []
    for part in iter_message_parts(m):
        pt = part.get("type")
        if pt == "text":
            t = (part.get("text") or "").strip()
            if t:
                blocks.append({"type": "text", "text": t})
        elif pt == "file":
            blocks.extend(file_part_to_openai_blocks(part))

    if not blocks:
        t = extract_text_from_message(m).strip()
        blocks = [{"type": "text", "text": t or "(empty message)"}]

    if len(blocks) == 1 and blocks[0]["type"] == "text":
        return {"role": "user", "content": blocks[0]["text"]}
    return {"role": "user", "content": blocks}


def raw_message_to_openai(m: dict) -> dict:
    role = m.get("role", "user")
    if role == "user":
        return build_user_openai_message(m)
    text = extract_text_from_message(m)
    return {"role": role, "content": text}


def mock_user_preview(last: dict | None) -> str:
    if not last or last.get("role") != "user":
        return "Hello"
    content = last.get("content")
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        bits: list[str] = []
        has_image = False
        for b in content:
            if not isinstance(b, dict):
                continue
            if b.get("type") == "text":
                bits.append(b.get("text", ""))
            elif b.get("type") == "image_url":
                has_image = True
        preview = " ".join(bits).strip()
        if has_image:
            preview = (preview + " [+ image attachment(s)]").strip()
        return preview or "[attachment(s) only]"
    return str(content)
