"""PPTX export endpoint — converts slides JSON to a downloadable PowerPoint file."""

from __future__ import annotations

from io import BytesIO
from typing import Any, List

from fastapi import APIRouter
from fastapi.responses import Response
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pydantic import BaseModel

router = APIRouter()


class ExportRequest(BaseModel):
    slides: List[dict[str, Any]]


@router.post("/api/export/pptx")
async def export_pptx(request: ExportRequest):
    """Generate a .pptx file from the slides array and return it as a download."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank slide

    for slide_data in request.slides:
        slide = prs.slides.add_slide(blank_layout)

        # Dark background matching the app's zinc-950 theme
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x09, 0x09, 0x0B)

        title_text = slide_data.get("title", "")
        body_items: list = slide_data.get("body", [])
        slide_type = slide_data.get("type", "concept")

        # Accent bar at the top (blue-to-purple gradient simulation via solid color)
        accent = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(0.12),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x60, 0x60, 0xFF)  # blue-500 approx
        accent.line.fill.background()

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(12.13), Inches(1.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title_text
        p = tf.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xF4, 0xF4, 0xF5)  # zinc-100

        # Slide type label (top-right)
        label_box = slide.shapes.add_textbox(
            Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.4)
        )
        label_tf = label_box.text_frame
        label_tf.text = slide_type.upper()
        if label_tf.paragraphs[0].runs:
            label_tf.paragraphs[0].runs[0].font.size = Pt(9)
            label_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x71, 0x71, 0x7A)  # zinc-500

        # Body bullets
        if body_items:
            body_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(2.1), Inches(12.13), Inches(4.8)
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for idx, bullet in enumerate(body_items):
                para = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
                para.text = f"  {bullet}"
                if para.runs:
                    para.runs[0].font.size = Pt(20)
                    para.runs[0].font.color.rgb = RGBColor(0xD4, 0xD4, 0xD8)  # zinc-300
                para.space_before = Pt(6)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="flowdeck-presentation.pptx"'},
    )
